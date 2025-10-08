import streamlit as st
import google.generativeai as genai
import fitz  # PyMuPDF
from pptx import Presentation
import json
import pandas as pd
import plotly.express as px
from utils.gemini_helper import GeminiHelper
from utils.template_checker import TemplateChecker
import os
import io


# Configure Streamlit page
st.set_page_config(
    page_title="DeckIQ - Pitch Deck Enhancer",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide Streamlit default UI elements
hide_streamlit_style = """
<style>
#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}
[data-testid="stToolbar"] {visibility: hidden !important;}
.stDeployButton {display: none;}
div[data-testid="stStatusWidget"] {display: none;}
#stDecoration {display: none;}
</style>
"""
st.markdown(hide_streamlit_style, unsafe_allow_html=True)


# Updated Gemini models (current as of October 2025)
GEMINI_MODELS = [
    "gemini-2.5-flash",        # Latest and fastest
    "gemini-2.0-flash",        # Stable fallback
    "gemini-1.5-flash-002",    # Legacy support
]


@st.cache_resource
def init_gemini():
    """Initialize Gemini with proper error handling and model selection"""
    api_key = st.secrets.get("GOOGLE_API_KEY") or os.getenv("GOOGLE_API_KEY")
    
    if not api_key:
        st.error("⚠️ Please set GOOGLE_API_KEY in Streamlit secrets or environment variables")
        st.info("Get your free API key from: https://aistudio.google.com/")
        return None, None

    try:
        genai.configure(api_key=api_key)

        # List available models
        try:
            available_models = [m.name.split('/')[-1] for m in genai.list_models()]
            st.sidebar.success("✅ Connected to Gemini API")
            st.sidebar.info(f"Available models: {len(available_models)}")
        except Exception as e:
            st.sidebar.warning(f"Could not list models: {str(e)}")
            available_models = []

        # Try each model until one works
        for model_name in GEMINI_MODELS:
            try:
                model = genai.GenerativeModel(model_name)
                # Test the model with a simple prompt
                test_response = model.generate_content(
                    "Respond with 'OK'",
                    generation_config=genai.GenerationConfig(
                        temperature=0.1,
                        max_output_tokens=10
                    )
                )
                if test_response.text:
                    st.sidebar.success(f"🤖 Active model: **{model_name}**")
                    return model, model_name
            except Exception as e:
                # Don't show error messages in sidebar - just try next model
                continue

        # If no models work, show error
        st.error("❌ Could not initialize any Gemini model. Please check your API key.")
        st.info("💡 Visit https://aistudio.google.com/ to verify your API key is active")
        return None, None

    except Exception as e:
        st.error(f"Error configuring Gemini: {str(e)}")
        return None, None


def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text
    except Exception as e:
        st.error(f"Error extracting PDF text: {str(e)}")
        return ""


def extract_text_from_pptx(file):
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error extracting PPTX text: {str(e)}")
        return ""


def show_api_setup_guide():
    """Show detailed API setup guide"""
    st.markdown("""
    ### 🔑 Google Gemini API Setup Guide

    **Step 1: Get Your Free API Key**
    1. Visit [**Google AI Studio**](https://aistudio.google.com/)
    2. Sign in with your Google account
    3. Click **"Get API Key"** or **"Create API Key"**
    4. Copy the generated key (starts with 'AIza...')

    **Step 2: Add to Streamlit**
    - Go to your app settings
    - Find "Secrets" section  
    - Add: `GOOGLE_API_KEY = "your_actual_key_here"`

    **Step 3: Verify Setup**
    - Refresh this page
    - You should see "✅ Connected to Gemini API" in the sidebar

    **Need Help?**
    - Make sure your API key is active
    - Check for any typos in the key
    - Ensure Gemini API is enabled in your Google Cloud project
    """)


def show_analysis_error(feature_name):
    """Display error message for failed analysis"""
    st.error(f"❌ Failed to generate {feature_name}. Please try again.")
    st.info("💡 **Troubleshooting tips:**\n- Check your internet connection\n- Verify API key is active\n- Try refreshing the page\n- Ensure your deck has readable text content")


# Main App
def main():
    st.title("📊 DeckIQ - Pitch Deck Enhancer")
    st.markdown("*Transform your pitch deck with AI-powered insights in seconds*")

    # Sidebar
    with st.sidebar:
        st.header("🚀 Features")
        st.markdown("""
        - **📑 Structure Analysis**: Investor-ready format
        - **🎤 Pitch Script**: Compelling 2-min presentation
        - **🎨 Design Tips**: Modern slide improvements
        - **📊 Benchmark Check**: Compare vs YC/Sequoia
        - **📄 One-Pager**: Executive summary
        """)
        
        st.markdown("---")
        
        # API status section
        st.markdown("**⚙️ API Status**")

    # Initialize Gemini
    model, model_name = init_gemini()
    if not model:
        show_api_setup_guide()
        return

    # File upload section
    st.markdown("### 📄 Upload Your Pitch Deck")
    uploaded_file = st.file_uploader(
        "Choose your pitch deck file",
        type=["pdf", "pptx"],
        help="Supported formats: PDF, PowerPoint (.pptx). Max size: 10MB"
    )

    if uploaded_file:
        # Extract text based on file type
        with st.spinner("🔍 Extracting content from your deck..."):
            if uploaded_file.type == "application/pdf":
                deck_text = extract_text_from_pdf(uploaded_file)
            else:  # pptx
                deck_text = extract_text_from_pptx(uploaded_file)

        if len(deck_text.strip()) < 50:
            st.warning("⚠️ Limited text detected. Ensure your deck contains readable text.")
            
            with st.expander("📋 See example format"):
                st.markdown("""
                **Good Format Example:**
                ```
                Problem: 73% of small businesses struggle with 24/7 support
                Solution: AI chatbot handling 90% of inquiries automatically
                Market: $15B customer service automation, growing 25% annually
                Traction: 150 customers, $45K MRR, 95% satisfaction
                ```
                """)
            return

        st.success(f"✅ Extracted {len(deck_text)} characters | {len(deck_text.split())} words")

        # Show preview
        with st.expander("👁️ Preview extracted content"):
            preview_text = deck_text[:800] + "..." if len(deck_text) > 800 else deck_text
            st.text_area("Content Preview", preview_text, height=200, disabled=True)

        # Initialize helpers
        helper = GeminiHelper(model, model_name)
        checker = TemplateChecker()

        # Analysis tabs
        tab1, tab2, tab3, tab4, tab5 = st.tabs([
            "📑 Structure", "🎤 Pitch Script", "🎨 Design", "📊 Benchmark", "📄 One-Pager"
        ])

        with tab1:
            st.markdown("### 📑 Structured Outline")
            st.markdown("Reorganize your deck into investor-ready sections")

            if st.button("🚀 Generate Structure", key="structure", type="primary"):
                with st.spinner("🤖 Analyzing deck structure..."):
                    try:
                        outline = helper.generate_structure(deck_text)
                        if outline and "Error" not in outline:
                            st.markdown("---")
                            st.markdown(outline)
                            st.download_button(
                                "📥 Download Structure",
                                outline,
                                file_name="pitch_structure.md",
                                mime="text/markdown"
                            )
                        else:
                            show_analysis_error("structure analysis")
                    except Exception as e:
                        show_analysis_error("structure analysis")

        with tab2:
            st.markdown("### 🎤 2-Minute Pitch Script")
            st.markdown("Transform your deck into a compelling presentation script")

            if st.button("🎯 Generate Script", key="script", type="primary"):
                with st.spinner("✍️ Crafting your pitch script..."):
                    try:
                        script = helper.generate_pitch_script(deck_text)
                        if script and "Error" not in script:
                            st.markdown("---")
                            st.markdown(script)
                            st.download_button(
                                "📥 Download Script",
                                script,
                                file_name="pitch_script.md",
                                mime="text/markdown"
                            )
                        else:
                            show_analysis_error("pitch script")
                    except Exception as e:
                        show_analysis_error("pitch script")

        with tab3:
            st.markdown("### 🎨 Design Suggestions")
            st.markdown("Get modern design recommendations for your slides")

            if st.button("🎨 Get Design Tips", key="design", type="primary"):
                with st.spinner("🎨 Analyzing design improvements..."):
                    try:
                        design_tips = helper.generate_design_suggestions(deck_text)
                        if design_tips and "Error" not in design_tips:
                            st.markdown("---")
                            st.markdown(design_tips)
                            st.download_button(
                                "📥 Download Design Guide",
                                design_tips,
                                file_name="design_suggestions.md",
                                mime="text/markdown"
                            )
                        else:
                            show_analysis_error("design suggestions")
                    except Exception as e:
                        show_analysis_error("design suggestions")

        with tab4:
            st.markdown("### 📊 Benchmark Analysis")
            st.markdown("Compare your deck against top-tier investor templates")

            template_choice = st.selectbox(
                "📋 Compare against:",
                ["Y Combinator", "Sequoia Capital"],
                help="Choose which template to benchmark against"
            )

            if st.button("📊 Run Benchmark", key="benchmark", type="primary"):
                with st.spinner("📈 Comparing against best practices..."):
                    try:
                        template_key = template_choice.lower().replace(" ", "_")
                        gaps = checker.check_template_gaps(deck_text, template_key)
                        benchmark_analysis = helper.generate_benchmark_analysis(
                            deck_text, gaps, template_choice
                        )

                        if benchmark_analysis and "Error" not in benchmark_analysis:
                            st.markdown("---")
                            
                            col1, col2 = st.columns([1, 1])

                            with col1:
                                st.subheader("✅ Coverage Score")
                                if gaps:
                                    missing_count = len(gaps)
                                    total_sections = len(
                                        checker.templates[template_key]['required_sections']
                                    )
                                    coverage = ((total_sections - missing_count) / total_sections) * 100
                                    st.metric("Coverage", f"{coverage:.0f}%")
                                    
                                    for gap in gaps:
                                        st.error(f"❌ Missing: **{gap}**")
                                else:
                                    st.success("🎉 All elements present!")
                                    st.metric("Coverage", "100%")

                            with col2:
                                st.subheader("📊 Requirements")
                                template_info = checker.templates[template_key]
                                st.info(f"**Required:** {len(template_info['required_sections'])}")
                                st.info(f"**Optional:** {len(template_info['optional_sections'])}")

                            st.markdown("---")
                            st.markdown(benchmark_analysis)
                            
                            st.download_button(
                                "📥 Download Report",
                                benchmark_analysis,
                                file_name=f"benchmark_{template_key}.md",
                                mime="text/markdown"
                            )
                        else:
                            show_analysis_error("benchmark analysis")
                    except Exception as e:
                        show_analysis_error("benchmark analysis")

        with tab5:
            st.markdown("### 📄 One-Page Executive Summary")
            st.markdown("Generate a concise investor-ready summary")

            if st.button("📝 Generate One-Pager", key="onepager", type="primary"):
                with st.spinner("📋 Creating executive summary..."):
                    try:
                        summary = helper.generate_one_pager(deck_text)
                        if summary and "Error" not in summary:
                            st.markdown("---")
                            st.markdown(summary)
                            st.download_button(
                                "📥 Download One-Pager",
                                summary,
                                file_name="executive_summary.md",
                                mime="text/markdown"
                            )
                        else:
                            show_analysis_error("one-pager")
                    except Exception as e:
                        show_analysis_error("one-pager")

    else:
        # Welcome screen
        st.markdown("## 👋 Welcome to DeckIQ")
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            st.markdown("""
            Transform your pitch deck with AI-powered analysis and improvements.

            ### 🌟 What you'll get:
            - **Structured outline** following investor best practices
            - **2-minute pitch script** for presentations
            - **Modern design suggestions** for visual appeal
            - **Benchmark analysis** vs top VC templates
            - **One-page summary** for investors
            """)
        
        with col2:
            st.info("""
            **📊 Supported formats:**
            - PDF files (.pdf)
            - PowerPoint (.pptx)
            
            **⚡ Processing:**
            - 10-30 seconds
            - Completely free
            """)

        with st.expander("📝 Try with sample content"):
            st.markdown("""
            **Sample Pitch to Test:**

            Create a PDF/PPT with this text:

            ---

            **Problem:** 73% of small businesses can't afford 24/7 support

            **Solution:** AI chatbot handling 90% of inquiries automatically

            **Market:** $15B automation market growing 25% annually

            **Traction:** 150 customers, $45K MRR, 95% satisfaction

            **Ask:** Raising $1.5M seed for engineering & sales expansion

            ---
            """)


if __name__ == "__main__":
    main()
